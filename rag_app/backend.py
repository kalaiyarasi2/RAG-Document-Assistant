# backend.py
import os
import pickle
import hashlib
from pathlib import Path
import numpy as np
import pandas as pd
from pptx import Presentation
from PyPDF2 import PdfReader
from docx import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from sentence_transformers import SentenceTransformer
import faiss
from groq import Groq


class RAGSystem:
    def __init__(self, raw_dir="raw_docs", processed_dir="processed", cache_dir="cache", api_key=None):
        self.raw_dir = Path(raw_dir)
        self.processed_dir = Path(processed_dir)
        self.cache_dir = Path(cache_dir)
        self.api_key = api_key or os.getenv("GROQ_API_KEY")
        if not self.api_key:
            raise ValueError("❌ No Groq API key provided. Set GROQ_API_KEY or pass api_key.")
        self.client = Groq(api_key=self.api_key)
        self.text_splitter = RecursiveCharacterTextSplitter(chunk_size=1024, chunk_overlap=100)
        self.embedder = SentenceTransformer("intfloat/e5-large-v2")
        self.index = None
        self.chunks = []
        for d in [self.raw_dir, self.processed_dir, self.cache_dir]:
            d.mkdir(exist_ok=True)

    def extract_text_from_pdf(self, path):
        try:
            reader = PdfReader(path)
            return "\n".join([(p.extract_text() or "") for p in reader.pages])
        except Exception as e:
            return f"[Error reading PDF: {e}]"

    def extract_text_from_docx(self, path):
        try:
            doc = Document(path)
            text = "\n".join(p.text for p in doc.paragraphs if p.text)
            for table in doc.tables:
                for row in table.rows:
                    text += "\n" + "\t".join(cell.text.strip() for cell in row.cells)
            return text
        except Exception as e:
            return f"[Error reading DOCX: {e}]"

    def extract_text_from_txt(self, path):
        try:
            with open(path, 'r', encoding='utf-8', errors='replace') as f:
                return f.read()
        except Exception as e:
            return f"[Error reading TXT: {e}]"

    def extract_text_from_csv(self, path):
        try:
            df = pd.read_csv(path)
            return df.to_string()
        except Exception as e:
            return f"[Error reading CSV: {e}]"

    def extract_text_from_excel(self, path):
        try:
            df = pd.read_excel(path)
            return df.to_string()
        except Exception as e:
            return f"[Error reading Excel: {e}]"

    def extract_text_from_pptx(self, path):
        try:
            prs = Presentation(path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            return f"[Error reading PPTX: {e}]"

    def load_documents(self):
        self.chunks = []
        extractors = {
            '.pdf': self.extract_text_from_pdf,
            '.docx': self.extract_text_from_docx,
            '.txt': self.extract_text_from_txt,
            '.csv': self.extract_text_from_csv,
            '.xlsx': self.extract_text_from_excel,
            '.pptx': self.extract_text_from_pptx
        }
        for file in self.raw_dir.iterdir():
            if file.suffix.lower() not in extractors or not file.is_file():
                continue
            try:
                print(f"📄 Processing {file.name}...")
                text = extractors[file.suffix.lower()](str(file))
                (self.processed_dir / f"{file.stem}.txt").write_text(text, encoding="utf-8")
                self.chunks.extend(self.text_splitter.split_text(text))
            except Exception as e:
                print(f"❌ Failed to process {file.name}: {e}")
        print(f"✅ Loaded {len(self.chunks)} chunks.")
        return self.chunks

    def calculate_hash(self):
        hash_md5 = hashlib.md5()
        files = sorted(self.raw_dir.glob("*.*"))
        for f in files:
            if f.is_file():
                with f.open("rb") as fb:
                    for chunk in iter(lambda: fb.read(4096), b""):
                        hash_md5.update(chunk)
        return hash_md5.hexdigest()

    def build_or_load_index(self, force=False):
        index_path = self.cache_dir / "faiss_index.bin"
        chunks_path = self.cache_dir / "chunks.pkl"
        hash_path = self.cache_dir / "doc_hash.txt"
        current_hash = self.calculate_hash()
        previous_hash = hash_path.read_text().strip() if hash_path.exists() else ""

        if not force and index_path.exists() and chunks_path.exists() and current_hash == previous_hash:
            print("🔁 Loading cached FAISS index...")
            self.index = faiss.read_index(str(index_path))
            with open(chunks_path, "rb") as f:
                self.chunks = pickle.load(f)
            return

        print("🆕 Building new FAISS index...")
        self.load_documents()
        if not self.chunks:
            raise ValueError("❌ No document content found. Upload valid files.")

        embeddings = self.embedder.encode(["passage: " + c for c in self.chunks], show_progress_bar=True)
        self.index = faiss.IndexFlatL2(embeddings.shape[1])
        self.index.add(np.array(embeddings))
        faiss.write_index(self.index, str(index_path))
        with open(chunks_path, "wb") as f:
            pickle.dump(self.chunks, f)
        hash_path.write_text(current_hash)
        print("💾 Index saved.")

    def query(self, question: str, model="llama3-8b-8192") -> tuple:
        if not self.index or not self.chunks:
            return "⚠️ Index not built. Please upload documents first.", ""

        try:
            query_emb = self.embedder.encode(["query: " + question])
            distances, indices = self.index.search(np.array(query_emb), k=6)
            context = "\n".join(
                [f"[Chunk {i+1}]\n{self.chunks[idx]}" for i, idx in enumerate(indices[0]) if idx < len(self.chunks)]
            )
            system_prompt = f"""Answer using only the provided context. If unsure, say 'I don't know'.
Context:
{context}"""
            response = self.client.chat.completions.create(
                model=model,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": question}
                ],
                temperature=0.3,
                max_tokens=512,
                top_p=0.9
            )
            return response.choices[0].message.content.strip(), context
        except Exception as e:
            return f"❌ API Error: {str(e)}", ""